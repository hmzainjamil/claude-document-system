#!/usr/bin/env python3
"""
Enterprise Document Factory — Hafiz Muhammad Zulqarnain
Zero formatting errors. Every format. Every time.

QA runner for PDF, DOCX, XLSX, PPTX documents.
Verifies file integrity, content presence, and page/slide counts.
Renders PDF page 1 to /tmp/doc_qa_page1.png as visual proof.

Usage:
    python3 ~/.claude/bin/doc-factory.py --qa ~/Downloads/report.pdf
    python3 ~/.claude/bin/doc-factory.py --qa ~/Downloads/report.docx
    python3 ~/.claude/bin/doc-factory.py --qa ~/Downloads/report.xlsx
    python3 ~/.claude/bin/doc-factory.py --qa ~/Downloads/report.pptx

    # With optional parameters:
    python3 ~/.claude/bin/doc-factory.py --qa ~/Downloads/report.xlsx --sheets Overview Campaigns
    python3 ~/.claude/bin/doc-factory.py --qa ~/Downloads/deck.pptx --slides 10

Exit codes:
    0 = PASS
    1 = FAIL (reason printed to stderr)
"""

import sys
import os
import argparse


# ── QA FUNCTIONS ─────────────────────────────────────────────────────────────

def qa_pdf(path: str) -> dict:
    """Full QA on a PDF using pymupdf."""
    try:
        import fitz
    except ImportError:
        raise RuntimeError("pymupdf not installed. Run: pip install pymupdf")

    result = {
        'format': 'PDF', 'path': path, 'pass': False,
        'reason': '', 'pages': 0, 'size_kb': 0,
    }

    if not os.path.exists(path):
        result['reason'] = f"FILE NOT FOUND: {path}"
        raise RuntimeError(result['reason'])

    size_bytes = os.path.getsize(path)
    result['size_kb'] = round(size_bytes / 1024, 1)

    if size_bytes < 10_240:
        result['reason'] = f"FILE TOO SMALL: {result['size_kb']}KB (min 10KB) — likely blank or corrupt"
        raise RuntimeError(result['reason'])

    doc = fitz.open(path)
    result['pages'] = len(doc)

    if result['pages'] == 0:
        result['reason'] = "ZERO PAGES — document is empty"
        raise RuntimeError(result['reason'])

    blank_pages = []
    thin_pages = []
    page_summaries = []

    for i, page in enumerate(doc):
        text = page.get_text().strip()
        char_count = len(text)
        preview = text[:60].replace('\n', ' ') if text else "(no text)"
        page_summaries.append(f"  Page {i+1}: {char_count} chars — {preview}")

        if char_count == 0:
            blank_pages.append(i + 1)
        elif char_count < 80:
            thin_pages.append((i + 1, char_count))

    if blank_pages:
        result['reason'] = f"BLANK PAGES DETECTED: pages {blank_pages} have 0 extractable characters"
        raise RuntimeError(result['reason'])

    if thin_pages:
        print(f"[PDF-QA] WARNING: thin pages (may be image-only cover): {thin_pages}", file=sys.stderr)

    # Render page 1 as PNG proof
    proof_path = "/tmp/doc_qa_page1.png"
    try:
        pix = doc[0].get_pixmap(dpi=96)
        pix.save(proof_path)
        print(f"[PDF-QA] Page 1 proof rendered → {proof_path}")
    except Exception as e:
        result['reason'] = f"PAGE RENDER FAILED: {e}"
        raise RuntimeError(result['reason'])

    doc.close()
    result['pass'] = True
    result['reason'] = "PASS"
    result['page_summaries'] = page_summaries
    return result


def qa_docx(path: str) -> dict:
    """Full QA on a DOCX file using python-docx."""
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx not installed. Run: pip install python-docx")

    result = {
        'format': 'DOCX', 'path': path, 'pass': False,
        'reason': '', 'paragraphs': 0, 'tables': 0, 'size_kb': 0,
    }

    if not os.path.exists(path):
        result['reason'] = f"FILE NOT FOUND: {path}"
        raise RuntimeError(result['reason'])

    size_bytes = os.path.getsize(path)
    result['size_kb'] = round(size_bytes / 1024, 1)

    if size_bytes < 5_120:
        result['reason'] = f"FILE TOO SMALL: {result['size_kb']}KB (min 5KB) — likely blank"
        raise RuntimeError(result['reason'])

    try:
        doc = Document(path)
    except Exception as e:
        result['reason'] = f"DOCX OPEN FAILED: {e}"
        raise RuntimeError(result['reason'])

    paras_with_text = [p for p in doc.paragraphs if p.text.strip()]
    result['paragraphs'] = len(paras_with_text)
    result['tables'] = len(doc.tables)

    if result['paragraphs'] == 0 and result['tables'] == 0:
        result['reason'] = "EMPTY DOCUMENT: zero text paragraphs and zero tables"
        raise RuntimeError(result['reason'])

    if result['paragraphs'] < 2 and result['tables'] == 0:
        result['reason'] = f"NEAR-EMPTY: only {result['paragraphs']} text paragraphs and no tables"
        raise RuntimeError(result['reason'])

    # Content summary
    result['content_summary'] = []
    for i, p in enumerate(paras_with_text[:5]):
        result['content_summary'].append(f"  Para {i+1}: {p.text[:80]}")
    if result['tables'] > 0:
        result['content_summary'].append(f"  Tables: {result['tables']} found")

    result['pass'] = True
    result['reason'] = "PASS"
    return result


def qa_xlsx(path: str, expected_sheets: list = None) -> dict:
    """Full QA on an XLSX file using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("openpyxl not installed. Run: pip install openpyxl")

    result = {
        'format': 'XLSX', 'path': path, 'pass': False,
        'reason': '', 'sheets': [], 'size_kb': 0,
    }

    if not os.path.exists(path):
        result['reason'] = f"FILE NOT FOUND: {path}"
        raise RuntimeError(result['reason'])

    size_bytes = os.path.getsize(path)
    result['size_kb'] = round(size_bytes / 1024, 1)

    if size_bytes < 3_072:
        result['reason'] = f"FILE TOO SMALL: {result['size_kb']}KB (min 3KB) — likely blank"
        raise RuntimeError(result['reason'])

    try:
        wb = load_workbook(path, read_only=True, data_only=True)
    except Exception as e:
        result['reason'] = f"XLSX OPEN FAILED: {e}"
        raise RuntimeError(result['reason'])

    result['sheets'] = wb.sheetnames

    if expected_sheets:
        missing = [s for s in expected_sheets if s not in result['sheets']]
        if missing:
            result['reason'] = f"MISSING SHEETS: {missing} — found: {result['sheets']}"
            raise RuntimeError(result['reason'])

    # Check each sheet has data
    thin_sheets = []
    sheet_summaries = []
    for sheet_name in result['sheets']:
        ws = wb[sheet_name]
        max_row = ws.max_row or 0
        max_col = ws.max_column or 0
        sheet_summaries.append(f"  Sheet '{sheet_name}': {max_row} rows x {max_col} cols")
        if max_row <= 1:
            thin_sheets.append(f"{sheet_name}(rows={max_row})")

    if thin_sheets and len(thin_sheets) == len(result['sheets']):
        result['reason'] = f"ALL SHEETS EMPTY OR HEADER-ONLY: {thin_sheets}"
        raise RuntimeError(result['reason'])

    wb.close()
    result['pass'] = True
    result['reason'] = "PASS"
    result['sheet_summaries'] = sheet_summaries
    return result


def qa_pptx(path: str, expected_slides: int = None) -> dict:
    """Full QA on a PPTX file using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx")

    result = {
        'format': 'PPTX', 'path': path, 'pass': False,
        'reason': '', 'slides': 0, 'size_kb': 0,
    }

    if not os.path.exists(path):
        result['reason'] = f"FILE NOT FOUND: {path}"
        raise RuntimeError(result['reason'])

    size_bytes = os.path.getsize(path)
    result['size_kb'] = round(size_bytes / 1024, 1)

    if size_bytes < 10_240:
        result['reason'] = f"FILE TOO SMALL: {result['size_kb']}KB (min 10KB) — likely blank"
        raise RuntimeError(result['reason'])

    try:
        prs = Presentation(path)
    except Exception as e:
        result['reason'] = f"PPTX OPEN FAILED: {e}"
        raise RuntimeError(result['reason'])

    result['slides'] = len(prs.slides)

    if result['slides'] == 0:
        result['reason'] = "ZERO SLIDES — presentation is empty"
        raise RuntimeError(result['reason'])

    if expected_slides is not None and result['slides'] != expected_slides:
        result['reason'] = f"SLIDE COUNT MISMATCH: got {result['slides']}, expected {expected_slides}"
        raise RuntimeError(result['reason'])

    blank_slides = []
    slide_summaries = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        first_text = texts[0][:60] if texts else "(no text)"
        slide_summaries.append(f"  Slide {i+1}: {first_text}")
        if not texts:
            blank_slides.append(i + 1)

    if blank_slides:
        result['reason'] = f"BLANK SLIDES DETECTED: slides {blank_slides} have no text content"
        raise RuntimeError(result['reason'])

    result['pass'] = True
    result['reason'] = "PASS"
    result['slide_summaries'] = slide_summaries
    return result


def run_document_qa(path: str, expected_sheets: list = None, expected_slides: int = None) -> dict:
    """
    Universal QA runner. Auto-detects format from extension.
    Raises RuntimeError if QA fails.
    Returns result dict with all details on success.
    """
    ext = os.path.splitext(path)[1].lower()

    print(f"\n{'='*60}")
    print(f"  DOCUMENT QA — Hafiz Muhammad Zulqarnain")
    print(f"  File  : {path}")
    print(f"  Format: {ext.upper()}")
    print(f"{'='*60}")

    if ext == '.pdf':
        result = qa_pdf(path)
    elif ext == '.docx':
        result = qa_docx(path)
    elif ext == '.xlsx':
        result = qa_xlsx(path, expected_sheets=expected_sheets)
    elif ext == '.pptx':
        result = qa_pptx(path, expected_slides=expected_slides)
    else:
        raise ValueError(f"Unsupported format: {ext} — supported: .pdf .docx .xlsx .pptx")

    # Print content summary
    if 'page_summaries' in result:
        print("\n[CONTENT SUMMARY — Pages]")
        for line in result['page_summaries']:
            print(line)
    if 'content_summary' in result:
        print("\n[CONTENT SUMMARY — DOCX]")
        for line in result['content_summary']:
            print(line)
    if 'sheet_summaries' in result:
        print("\n[CONTENT SUMMARY — Sheets]")
        for line in result['sheet_summaries']:
            print(line)
    if 'slide_summaries' in result:
        print("\n[CONTENT SUMMARY — Slides]")
        for line in result['slide_summaries']:
            print(line)

    print(f"\n{'='*60}")
    print(f"  STATUS : {result['reason']}")
    print(f"  SIZE   : {result['size_kb']} KB")
    if 'pages' in result:
        print(f"  PAGES  : {result['pages']}")
    if 'slides' in result:
        print(f"  SLIDES : {result['slides']}")
    if 'sheets' in result:
        print(f"  SHEETS : {result['sheets']}")
    if 'paragraphs' in result:
        print(f"  PARAS  : {result['paragraphs']} | TABLES: {result['tables']}")
    print(f"{'='*60}\n")

    return result


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Enterprise Document QA — Hafiz Muhammad Zulqarnain",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python3 doc-factory.py --qa ~/Downloads/report.pdf
  python3 doc-factory.py --qa ~/Downloads/report.docx
  python3 doc-factory.py --qa ~/Downloads/data.xlsx --sheets Overview Campaigns
  python3 doc-factory.py --qa ~/Downloads/deck.pptx --slides 10

Exit codes: 0=PASS  1=FAIL
        """
    )
    parser.add_argument('--qa', metavar='FILE', help='Path to document to QA')
    parser.add_argument('--sheets', nargs='+', metavar='SHEET', help='Expected sheet names (XLSX only)')
    parser.add_argument('--slides', type=int, metavar='N', help='Expected slide count (PPTX only)')

    args = parser.parse_args()

    if not args.qa:
        parser.print_help()
        sys.exit(0)

    path = os.path.expanduser(args.qa)

    try:
        result = run_document_qa(
            path,
            expected_sheets=args.sheets,
            expected_slides=args.slides,
        )
        print(f"✓ QA PASSED: {path}")
        sys.exit(0)
    except RuntimeError as e:
        print(f"\n✗ QA FAILED: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"\n✗ UNEXPECTED ERROR: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
